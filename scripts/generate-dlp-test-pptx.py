#!/usr/bin/env python3
"""Generate a PowerPoint deck summarizing the Agent Warden DLP test results."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colours ──
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_BLUE = RGBColor(0x00, 0x33, 0x66)
MID_BLUE = RGBColor(0x00, 0x78, 0xD4)  # Microsoft blue
LIGHT_BLUE = RGBColor(0xDE, 0xEC, 0xF9)
GREEN = RGBColor(0x10, 0x7C, 0x10)
RED = RGBColor(0xD1, 0x34, 0x38)
GRAY = RGBColor(0x60, 0x60, 0x60)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
ACCENT_PURPLE = RGBColor(0x5B, 0x2D, 0x8E)
CODE_BG = RGBColor(0x1E, 0x1E, 0x2E)   # Dark code background
CODE_FG = RGBColor(0xCD, 0xD6, 0xF4)   # Light code text
CODE_KW = RGBColor(0x89, 0xB4, 0xFA)   # Blue keywords
CODE_STR = RGBColor(0xA6, 0xE3, 0xA1)  # Green strings
CODE_CMT = RGBColor(0x6C, 0x70, 0x86)  # Gray comments

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


# ── Helpers ──
def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=BLACK, align=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_rich_textbox(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.text_frame.word_wrap = True
    return txBox.text_frame


def add_para(tf, text, font_size=16, bold=False, color=BLACK, align=PP_ALIGN.LEFT,
             space_before=Pt(4), space_after=Pt(2), font_name="Segoe UI"):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_table_row_cells(table, row_idx, values, font_size=12, bold=False, color=BLACK, fill=None):
    for col_idx, val in enumerate(values):
        cell = table.cell(row_idx, col_idx)
        cell.text = str(val)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.bold = bold
            paragraph.font.color.rgb = color
            paragraph.font.name = "Segoe UI"
        if fill:
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill


def add_stripe_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_code_box(slide, left, top, width, height, lines, font_size=9):
    """Add a dark-themed code block. lines = list of (text, color) tuples or strings."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.color.rgb = RGBColor(0x40, 0x40, 0x50)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(0)
        p.space_after = Pt(0)

        if isinstance(line, str):
            run = p.add_run()
            run.text = line
            run.font.size = Pt(font_size)
            run.font.color.rgb = CODE_FG
            run.font.name = "Menlo"
        elif isinstance(line, list):
            # list of (text, color) tuples for syntax highlighting
            for text, color in line:
                run = p.add_run()
                run.text = text
                run.font.size = Pt(font_size)
                run.font.color.rgb = color
                run.font.name = "Menlo"
        else:
            text, color = line
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = "Menlo"
    return shape


# ════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, WHITE)
add_stripe_rect(slide, 0, 0, 13.333, 0.15, MID_BLUE)
add_stripe_rect(slide, 0, 7.35, 13.333, 0.15, MID_BLUE)

add_textbox(slide, 1.5, 1.8, 10.3, 1.2,
            "Agent Warden — Purview DLP Integration",
            font_size=36, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)

add_textbox(slide, 1.5, 3.2, 10.3, 0.8,
            "End-to-End Test Results",
            font_size=28, color=MID_BLUE, align=PP_ALIGN.CENTER)

add_textbox(slide, 1.5, 4.4, 10.3, 0.6,
            "Microsoft Purview DLP × Microsoft Graph processContent API",
            font_size=18, color=GRAY, align=PP_ALIGN.CENTER)

add_textbox(slide, 1.5, 5.6, 10.3, 0.5,
            "Tenant: dab94ed2 (aprforazure)  •  Date: April 15, 2026  •  5/5 BLOCK ✓",
            font_size=16, color=GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Architecture
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_stripe_rect(slide, 0, 0, 13.333, 0.08, MID_BLUE)

add_textbox(slide, 0.5, 0.3, 12, 0.7,
            "Architecture — Cross-Tenant DLP Enforcement",
            font_size=28, bold=True, color=DARK_BLUE)

# Left column — Flow diagram as text boxes
def add_box(slide, left, top, width, height, text, fill_color, text_color=WHITE, font_size=12, bold=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = text_color
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return shape

def add_arrow(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = MID_BLUE
    shape.line.fill.background()
    return shape

# AKS Tenant (Contoso)
add_box(slide, 0.5, 1.3, 5.5, 0.55, "AKS Cluster — aks-agentwarden-dev (Contoso Tenant 9a72f9b7)", DARK_BLUE, font_size=13)

add_box(slide, 0.8, 2.1, 2.2, 0.45, "OpenClaw Gateway", MID_BLUE, font_size=11)
add_box(slide, 3.3, 2.1, 2.4, 0.45, "DLP Plugin (v0.5.5)", ACCENT_PURPLE, font_size=11)

add_box(slide, 0.8, 2.75, 4.9, 0.45, "L1: Prompt Guard  →  L2: Output Scanner  →  L2b: Response Scanner  →  L3: Input Audit", RGBColor(0x40, 0x40, 0x40), font_size=9)

add_arrow(slide, 3.0, 3.4, 0.35, 0.45)

# Purview Tenant
add_box(slide, 0.5, 4.05, 5.5, 0.55, "Purview Tenant — dab94ed2 (aprforazure.onmicrosoft.com)", RGBColor(0x00, 0x6B, 0x6B), font_size=13)

add_box(slide, 0.8, 4.8, 2.5, 0.45, "ClientSecretCredential", RGBColor(0x00, 0x89, 0x89), font_size=10)
add_box(slide, 3.6, 4.8, 2.2, 0.45, "Graph processContent", RGBColor(0x00, 0x89, 0x89), font_size=10)

add_box(slide, 0.8, 5.5, 5.0, 0.45, "App: d94c93dd  •  User: Frank (21bbd518)  •  E5 Licensed", RGBColor(0x40, 0x40, 0x40), font_size=10)

# Right column — Key details
tf = add_rich_textbox(slide, 6.8, 1.3, 6.0, 5.5)
add_para(tf, "Cross-Tenant Auth Flow", font_size=18, bold=True, color=DARK_BLUE)
add_para(tf, "1. DLP Plugin acquires token from Purview tenant (dab94ed2) using ClientSecretCredential", font_size=13, color=BLACK)
add_para(tf, "2. Calls POST /users/{userId}/dataSecurityAndGovernance/processContent", font_size=13, color=BLACK)
add_para(tf, "3. Graph evaluates content against DLP policies and returns policyActions", font_size=13, color=BLACK)
add_para(tf, "4. Plugin enforces: restrictAccess → block content, empty → allow", font_size=13, color=BLACK)

add_para(tf, "", font_size=8)
add_para(tf, "DLP Layers", font_size=18, bold=True, color=DARK_BLUE)
add_para(tf, "• L1 Prompt Guard — injects DLP policy into LLM system context", font_size=13, color=BLACK)
add_para(tf, "• L2 Output Scanner — scans tool results via processContent (sync)", font_size=13, color=BLACK)
add_para(tf, "• L2b Response Scanner — scans LLM response before delivery", font_size=13, color=BLACK)
add_para(tf, "• L3 Input Audit — scans user input, taints thread on violation", font_size=13, color=BLACK)

add_para(tf, "", font_size=8)
add_para(tf, "Key Graph Permissions", font_size=18, bold=True, color=DARK_BLUE)
add_para(tf, "• Content.Process.All  •  ProtectionScopes.Compute.All", font_size=13, color=BLACK)
add_para(tf, "• InformationProtectionPolicy.Read.All  •  ContentActivity.Write", font_size=13, color=BLACK)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 3 — DLP Policy Configuration
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_stripe_rect(slide, 0, 0, 13.333, 0.08, MID_BLUE)

add_textbox(slide, 0.5, 0.3, 12, 0.7,
            "DLP Policy Configuration",
            font_size=28, bold=True, color=DARK_BLUE)

# Policy table
cols = 2
rows = 9
table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(5.8), Inches(4.5))
table = table_shape.table
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(3.8)

policy_data = [
    ("Property", "Value"),
    ("Policy Name", "Agent Warden - Entra DLP"),
    ("GUID", "1cb19044-daf9-4bc8-af62-db7b032dd43d"),
    ("Workload", "Applications"),
    ("EnforcementPlanes", "Entra"),
    ("LocationSource", "Entra"),
    ("LocationType", "Individual"),
    ("App ID", "d94c93dd-3c80-4f3d-9671-8b71a7dccafa"),
    ("Mode", "Enable"),
]

for i, (k, v) in enumerate(policy_data):
    if i == 0:
        add_table_row_cells(table, i, [k, v], font_size=12, bold=True, color=WHITE, fill=DARK_BLUE)
    else:
        fill = LIGHT_GRAY if i % 2 == 0 else WHITE
        add_table_row_cells(table, i, [k, v], font_size=12, fill=fill)

# Rule table
add_textbox(slide, 0.5, 5.9, 5.8, 0.5,
            "Rule: Block PII via Entra App  •  Action: RestrictAccess (UploadText: Block)",
            font_size=13, bold=True, color=DARK_BLUE)

# SIT categories (right side)
tf = add_rich_textbox(slide, 6.8, 1.2, 6.0, 5.5)
add_para(tf, "28 Sensitive Information Types (SITs)", font_size=18, bold=True, color=DARK_BLUE)

add_para(tf, "", font_size=6)
add_para(tf, "Credit Card", font_size=14, bold=True, color=ACCENT_PURPLE)
add_para(tf, "Credit Card Number (Visa, MC, AmEx, Discover, JCB, UnionPay, Diners)", font_size=12, color=BLACK)

add_para(tf, "", font_size=4)
add_para(tf, "U.S. PII", font_size=14, bold=True, color=ACCENT_PURPLE)
add_para(tf, "U.S. Social Security Number (SSN), U.S. Physical Addresses", font_size=12, color=BLACK)

add_para(tf, "", font_size=4)
add_para(tf, "Passports (16 countries)", font_size=14, bold=True, color=ACCENT_PURPLE)
add_para(tf, "U.S./U.K., Canada, Australia, Japan, France, Germany, Italy, Spain, Netherlands, Belgium, Sweden, Finland, Austria, Ireland, South Korea, Taiwan, Poland", font_size=12, color=BLACK)

add_para(tf, "", font_size=4)
add_para(tf, "Russia", font_size=14, bold=True, color=ACCENT_PURPLE)
add_para(tf, "Passport (Domestic + International), Taxpayer ID (INN), Physical Addresses", font_size=12, color=BLACK)

add_para(tf, "", font_size=4)
add_para(tf, "Japan", font_size=14, bold=True, color=ACCENT_PURPLE)
add_para(tf, "My Number (Personal + Corporate), Resident Registration, Physical Addresses", font_size=12, color=BLACK)

add_para(tf, "", font_size=6)
add_para(tf, "Root Cause Fix (2026-04-15)", font_size=16, bold=True, color=RED)
add_para(tf, 'The -Locations JSON must include LocationSource:"Entra" and LocationType:"Individual" for custom Entra-registered app IDs. Without these fields, the DLP system cannot resolve the app ID → "Location is invalid".', font_size=11, color=GRAY)
add_para(tf, "Applications workload requires a separate policy (cannot combine with Exchange/SPO/OD4B).", font_size=11, color=GRAY)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Test Scenarios
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_stripe_rect(slide, 0, 0, 13.333, 0.08, MID_BLUE)

add_textbox(slide, 0.5, 0.3, 12, 0.7,
            "Test Scenarios — Checksum-Valid Fake PII Data",
            font_size=28, bold=True, color=DARK_BLUE)

add_textbox(slide, 0.5, 1.0, 12, 0.4,
            "All fake data passes checksum/format validation (Luhn, check digits). Keywords in correct language per country.",
            font_size=14, color=GRAY)

# Scenario table
cols = 5
rows = 6
table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.0))
table = table_shape.table
table.columns[0].width = Inches(0.5)
table.columns[1].width = Inches(2.3)
table.columns[2].width = Inches(3.0)
table.columns[3].width = Inches(3.5)
table.columns[4].width = Inches(3.0)

headers = ["#", "Scenario", "Key PII Data", "SITs Triggered", "Notes"]
add_table_row_cells(table, 0, headers, font_size=12, bold=True, color=WHITE, fill=DARK_BLUE)

scenarios = [
    ("1", "China PII\n(ChatGPT input)",
     "UnionPay: 6225887632109878\nPassport: G12345678\nID: 11010519900307783X",
     "Credit Card Number",
     "UnionPay Luhn-valid\nChinese + English keywords"),
    ("2", "US PII\n(Gemini file)",
     "Visa: 4242424242424242\nSSN: 234-67-8901\nPassport: 285194736",
     "Credit Card Number\nU.S. SSN\nU.S./U.K. Passport",
     "All checksum-valid\nEnglish keywords"),
    ("3", "Japan PII\n(ZIP attachment)",
     "JCB: 3528234567890126\nMy Number: 123456789018\nPassport: TR1234567",
     "Credit Card Number\nJapanese My Number Personal\nJapan Passport",
     "JCB Luhn-valid\nJapanese keywords (パスポート, マイナンバー)"),
    ("4", "France PII\n(Claude.ai)",
     "MasterCard: 5555555555554444\nINSÉE: 185067283009381\nPassport: 10AB12345",
     "Credit Card Number\nFrance Passport",
     "MC Luhn-valid, INSÉE key=81\nFrench keywords"),
    ("5", "Russia PII\n(photo text)",
     "MC 2221: 2221001234567896\nINN: 770123456703\nPassport: 45 1234567",
     "Credit Card Number\nRussian Passport (Domestic)\nRussian Taxpayer ID",
     "MC 2221 prefix (not Mir 2200)\nRussian keywords (Номер паспорта)"),
]

for i, row in enumerate(scenarios):
    fill = LIGHT_BLUE if i % 2 == 0 else WHITE
    add_table_row_cells(table, i + 1, row, font_size=10, fill=fill)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Test Results
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_stripe_rect(slide, 0, 0, 13.333, 0.08, MID_BLUE)

add_textbox(slide, 0.5, 0.3, 12, 0.7,
            "Test Results — 5/5 BLOCK ✓",
            font_size=28, bold=True, color=DARK_BLUE)

# Config box
tf = add_rich_textbox(slide, 0.5, 1.2, 4.5, 1.5)
add_para(tf, "Test Configuration", font_size=16, bold=True, color=DARK_BLUE)
add_para(tf, "Tenant:  dab94ed2 (aprforazure)", font_size=12, color=BLACK)
add_para(tf, "User:     Frank (21bbd518) — E5 Licensed", font_size=12, color=BLACK)
add_para(tf, "App:      d94c93dd (Agent Warden Purview DLP)", font_size=12, color=BLACK)
add_para(tf, "API:      Graph v1.0 processContent", font_size=12, color=BLACK)
add_para(tf, "Date:     April 15, 2026", font_size=12, color=BLACK)

# Results table
cols = 5
rows = 6
table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(3.0), Inches(8.0), Inches(3.5))
table = table_shape.table
table.columns[0].width = Inches(0.5)
table.columns[1].width = Inches(3.0)
table.columns[2].width = Inches(1.0)
table.columns[3].width = Inches(1.5)
table.columns[4].width = Inches(2.0)

headers = ["#", "Scenario", "HTTP", "Result", "Action"]
add_table_row_cells(table, 0, headers, font_size=13, bold=True, color=WHITE, fill=DARK_BLUE)

results = [
    ("1", "China PII (valid Luhn)", "200", "BLOCK ✓", "restrictAccess → block"),
    ("2", "US PII (valid Luhn + SSN)", "200", "BLOCK ✓", "restrictAccess → block"),
    ("3", "Japan PII (valid Luhn + My Number)", "200", "BLOCK ✓", "restrictAccess → block"),
    ("4", "France PII (valid Luhn + INSÉE)", "200", "BLOCK ✓", "restrictAccess → block"),
    ("5", "Russia PII (valid Luhn + INN)", "200", "BLOCK ✓", "restrictAccess → block"),
]

for i, row in enumerate(results):
    fill = RGBColor(0xE8, 0xF5, 0xE9) if i % 2 == 0 else WHITE
    add_table_row_cells(table, i + 1, row, font_size=12, fill=fill)
    # Color the Result column green
    cell = table.cell(i + 1, 3)
    for p in cell.text_frame.paragraphs:
        p.font.color.rgb = GREEN
        p.font.bold = True

# Summary box
add_box(slide, 5.5, 1.3, 3.5, 1.2,
        "5 / 5 PASSED\nAll Scenarios BLOCK",
        GREEN, WHITE, font_size=22, bold=True)

# Key takeaway
tf = add_rich_textbox(slide, 9.5, 1.2, 3.5, 5.5)
add_para(tf, "Key Takeaways", font_size=18, bold=True, color=DARK_BLUE)
add_para(tf, "", font_size=6)
add_para(tf, "✓  Entra DLP policy fully distributed and enforcing on tenant dab94ed2", font_size=12, color=GREEN)
add_para(tf, "", font_size=4)
add_para(tf, "✓  All 5 country scenarios blocked via restrictAccess action", font_size=12, color=GREEN)
add_para(tf, "", font_size=4)
add_para(tf, "✓  28 SITs covering credit cards, passports, SSN, addresses, and national IDs", font_size=12, color=GREEN)
add_para(tf, "", font_size=4)
add_para(tf, "✓  Cross-tenant auth working (Contoso AKS → aprforazure Purview)", font_size=12, color=GREEN)
add_para(tf, "", font_size=8)
add_para(tf, "Next Steps", font_size=16, bold=True, color=DARK_BLUE)
add_para(tf, "• Switch Helm values to dab94ed2 + Frank", font_size=12, color=BLACK)
add_para(tf, "• Helm upgrade + pod restart", font_size=12, color=BLACK)
add_para(tf, "• E2E Telegram test with live DLP", font_size=12, color=BLACK)
add_para(tf, "• Retire old tenant 2cf24558", font_size=12, color=BLACK)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Code: processContent API Call
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_stripe_rect(slide, 0, 0, 13.333, 0.08, MID_BLUE)

add_textbox(slide, 0.5, 0.3, 12, 0.7,
            "Code — processContent API Request Body",
            font_size=28, bold=True, color=DARK_BLUE)

add_textbox(slide, 0.5, 0.9, 12, 0.4,
            "POST /v1.0/users/{userId}/dataSecurityAndGovernance/processContent  •  Source: purview-client.ts",
            font_size=13, color=GRAY)

add_code_box(slide, 0.5, 1.5, 6.0, 5.5, [
    [("// processContent request body", CODE_CMT)],
    [("const ", CODE_KW), ("body = {", CODE_FG)],
    [("  contentToProcess: {", CODE_FG)],
    [("    contentEntries: [{", CODE_FG)],
    [('      "@odata.type"', CODE_STR), (": ", CODE_FG), ('"microsoft.graph.processConversationMetadata"', CODE_STR), (",", CODE_FG)],
    [("      identifier: ", CODE_FG), ("crypto.randomUUID()", CODE_KW), (",", CODE_FG)],
    [("      content: {", CODE_FG)],
    [('        "@odata.type"', CODE_STR), (": ", CODE_FG), ('"microsoft.graph.textContent"', CODE_STR), (",", CODE_FG)],
    [("        data: ", CODE_FG), ("text", CODE_KW), ("  ", CODE_FG), ("// ← content to scan", CODE_CMT)],
    [("      },", CODE_FG)],
    [("      correlationId: ", CODE_FG), ("crypto.randomUUID()", CODE_KW), (",", CODE_FG)],
    [("      sequenceNumber: ", CODE_FG), ("0", CODE_STR), (",", CODE_FG)],
    [("      isTruncated: ", CODE_FG), ("false", CODE_KW)],
    [("    }],", CODE_FG)],
    [("    activityMetadata: { activity: ", CODE_FG), ('"uploadText"', CODE_STR), (" },", CODE_FG)],
    [("    protectedAppMetadata: {", CODE_FG)],
    [('      name: ', CODE_FG), ('"Agent Warden"', CODE_STR), (",", CODE_FG)],
    [("      applicationLocation: {", CODE_FG)],
    [('        "@odata.type"', CODE_STR), (": ", CODE_FG)],
    [('          "#microsoft.graph.policyLocationApplication"', CODE_STR), (",", CODE_FG)],
    [("        value: ", CODE_FG), ("CLIENT_ID", CODE_KW)],
    [("      }", CODE_FG)],
    [("    }", CODE_FG)],
    [("  }", CODE_FG)],
    [("};", CODE_FG)],
], font_size=9)

# Response parsing code
add_code_box(slide, 6.8, 1.5, 6.0, 3.0, [
    [("// Parse processContent response", CODE_CMT)],
    [("const ", CODE_KW), ("data = ", CODE_FG), ("await ", CODE_KW), ("resp.json();", CODE_FG)],
    [("const ", CODE_KW), ("actions = data.policyActions ?? [];", CODE_FG)],
    "",
    [("// Check for restrictAccess or block actions", CODE_CMT)],
    [("const ", CODE_KW), ("blocked = actions.some(", CODE_FG)],
    [("  (a) => a.action === ", CODE_FG), ('"restrictAccess"', CODE_STR)],
    [("      || a.action === ", CODE_FG), ('"block"', CODE_STR)],
    [("  );", CODE_FG)],
    "",
    [("// Example response when BLOCKED:", CODE_CMT)],
    [("// {", CODE_CMT)],
    [('//   "policyActions": [{', CODE_CMT)],
    [('//     "@odata.type":  ', CODE_CMT)],
    [('//       "#microsoft.graph.restrictAccessAction",', CODE_CMT)],
    [('//     "action": "restrictAccess",', CODE_CMT)],
    [('//     "restrictionAction": "block"', CODE_CMT)],
    [("//   }]", CODE_CMT)],
    [("// }", CODE_CMT)],
], font_size=9)

# Auth code
add_code_box(slide, 6.8, 4.7, 6.0, 2.3, [
    [("// Cross-tenant authentication (purview-client.ts)", CODE_CMT)],
    [("const ", CODE_KW), ("credential = ", CODE_FG), ("new ", CODE_KW), ("ClientSecretCredential(", CODE_FG)],
    [("  ", CODE_FG), ("PURVIEW_TENANT_ID", CODE_KW), (",  ", CODE_FG), ("// dab94ed2-...", CODE_CMT)],
    [("  ", CODE_FG), ("PURVIEW_CLIENT_ID", CODE_KW), (",  ", CODE_FG), ("// d94c93dd-...", CODE_CMT)],
    [("  ", CODE_FG), ("PURVIEW_CLIENT_SECRET", CODE_KW)],
    [(");", CODE_FG)],
    "",
    [("const ", CODE_KW), ("token = ", CODE_FG), ("await ", CODE_KW), ("credential.getToken(", CODE_FG)],
    [("  ", CODE_FG), ('"https://graph.microsoft.com/.default"', CODE_STR)],
    [(");", CODE_FG)],
], font_size=9)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Code: DLP Plugin Hook Architecture
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_stripe_rect(slide, 0, 0, 13.333, 0.08, MID_BLUE)

add_textbox(slide, 0.5, 0.3, 12, 0.7,
            "Code — DLP Plugin Hook Registration (index.ts)",
            font_size=28, bold=True, color=DARK_BLUE)

add_textbox(slide, 0.5, 0.9, 12, 0.4,
            "Four layers registered as OpenClaw plugin hooks  •  Source: agent-warden-purview-dlp/src/index.ts",
            font_size=13, color=GRAY)

# L1 + L3
add_code_box(slide, 0.5, 1.5, 6.0, 2.8, [
    [("// L1: Prompt Guard — inject DLP policy into LLM context", CODE_CMT)],
    [("api.hooks.register(", CODE_FG), ('"before_agent_start"', CODE_STR), (", ", CODE_FG), ("async ", CODE_KW), ("() => {", CODE_FG)],
    [("  ", CODE_FG), ("return ", CODE_KW), ("{ prependContext: DLP_SYSTEM_PROMPT };", CODE_FG)],
    [("});", CODE_FG)],
    "",
    [("// L3: Input Audit — scan inbound user messages", CODE_CMT)],
    [("api.hooks.register(", CODE_FG), ('"message_received"', CODE_STR), (", ", CODE_FG), ("async ", CODE_KW), ("(ev) => {", CODE_FG)],
    [("  ", CODE_FG), ("const ", CODE_KW), ("result = ", CODE_FG), ("await ", CODE_KW), ("purview.processContent(", CODE_FG)],
    [("    ev.content, ", CODE_FG), ('"uploadText"', CODE_STR), (");", CODE_FG)],
    [("  ", CODE_FG), ("if ", CODE_KW), ("(!result.allowed) tracker.taint(threadId);", CODE_FG)],
    [("  ", CODE_FG), ("// void hook — cannot block, logs + taints", CODE_CMT)],
    [("});", CODE_FG)],
], font_size=9)

# L2
add_code_box(slide, 0.5, 4.5, 6.0, 2.7, [
    [("// L2: Output Scanner — scan tool results (SYNC)", CODE_CMT)],
    [("api.hooks.register(", CODE_FG), ('"tool_result_persist"', CODE_STR), (", (ev) => {", CODE_FG)],
    [("  ", CODE_FG), ("const ", CODE_KW), ("mode = purview.getExecutionMode(", CODE_FG), ('"uploadText"', CODE_STR), (");", CODE_FG)],
    [("  ", CODE_FG), ("if ", CODE_KW), ('(mode === ', CODE_FG), ('"evaluateInline"', CODE_STR), (") {", CODE_FG)],
    [("    ", CODE_FG), ("// Synchronous scan via spawnSync + curl", CODE_CMT)],
    [("    ", CODE_FG), ("const ", CODE_KW), ("r = purview.processContentSync(content);", CODE_FG)],
    [("    ", CODE_FG), ("if ", CODE_KW), ("(!r.allowed) {", CODE_FG)],
    [("      tracker.taint(threadId);", CODE_FG)],
    [("      ", CODE_FG), ("return ", CODE_KW), ('{ message: "[DLP] Redacted" };', CODE_FG)],
    [("    }", CODE_FG)],
    [("  }", CODE_FG)],
    [("});", CODE_FG)],
], font_size=9)

# L2b
add_code_box(slide, 6.8, 1.5, 6.0, 2.8, [
    [("// L2b: Response Scanner — last line of defense", CODE_CMT)],
    [("// Only in enforce mode, requires streaming OFF", CODE_CMT)],
    [("api.hooks.register(", CODE_FG), ('"message_sending"', CODE_STR), (", ", CODE_FG), ("async ", CODE_KW), ("(ev) => {", CODE_FG)],
    [("  ", CODE_FG), ("// Block unconditionally if thread was tainted", CODE_CMT)],
    [("  ", CODE_FG), ("if ", CODE_KW), ("(tracker.isTainted(threadId)) {", CODE_FG)],
    [("    tracker.clearTaint(threadId);", CODE_FG)],
    [("    ", CODE_FG), ("return ", CODE_KW), ('{ content: "[DLP] Blocked" };', CODE_FG)],
    [("  }", CODE_FG)],
    [("  ", CODE_FG), ("const ", CODE_KW), ("r = ", CODE_FG), ("await ", CODE_KW), ("purview.processContent(", CODE_FG)],
    [("    ev.content, ", CODE_FG), ('"uploadText"', CODE_STR), (");", CODE_FG)],
    [("  ", CODE_FG), ("if ", CODE_KW), ('(!r.allowed) ', CODE_FG), ("return ", CODE_KW), ('{ content: "[DLP]..." };', CODE_FG)],
    [("});", CODE_FG)],
], font_size=9)

# DLP policy creation
add_code_box(slide, 6.8, 4.5, 6.0, 2.7, [
    [("# PowerShell — Create Entra DLP Policy", CODE_CMT)],
    [("$locations", CODE_KW), (" = ", CODE_FG), ("'[{", CODE_STR)],
    [('  "Workload": "Applications",', CODE_STR)],
    [('  "Location": "d94c93dd-...",', CODE_STR)],
    [('  "LocationSource": "Entra",', CODE_STR), ("   # REQUIRED!", CODE_CMT)],
    [('  "LocationType": "Individual",', CODE_STR), (" # REQUIRED!", CODE_CMT)],
    [("  ...}]'", CODE_STR)],
    "",
    [("New-DlpCompliancePolicy", CODE_KW), (' -Name "Agent Warden"', CODE_FG)],
    [("  -Locations ", CODE_FG), ("$locations", CODE_KW)],
    [("  -EnforcementPlanes ", CODE_FG), ('@("Entra")', CODE_STR)],
], font_size=9)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Test Code: scanText Function
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_stripe_rect(slide, 0, 0, 13.333, 0.08, MID_BLUE)

add_textbox(slide, 0.5, 0.3, 12, 0.7,
            "Test Code — scanText() & Test Runner",
            font_size=28, bold=True, color=DARK_BLUE)

add_textbox(slide, 0.5, 0.9, 12, 0.4,
            "Source: agent-warden-purview-dlp/test/test-dlp-scenarios-valid.ts",
            font_size=13, color=GRAY)

# scanText function
add_code_box(slide, 0.5, 1.5, 6.0, 5.5, [
    [("async function ", CODE_KW), ("scanText", CODE_FG), ("(text: ", CODE_KW), ("string", CODE_STR), ("): Promise<ScanResult> {", CODE_FG)],
    [("  ", CODE_FG), ("const ", CODE_KW), ("token = ", CODE_FG), ("await ", CODE_KW), ("getToken();", CODE_FG)],
    [("  ", CODE_FG), ("const ", CODE_KW), ("body = {", CODE_FG)],
    [("    contentToProcess: {", CODE_FG)],
    [("      contentEntries: [{", CODE_FG)],
    [('        "@odata.type": ', CODE_FG), ('"...processConversationMetadata"', CODE_STR), (",", CODE_FG)],
    [("        identifier: crypto.randomUUID(),", CODE_FG)],
    [("        content: { data: text },", CODE_FG)],
    [("        correlationId: crypto.randomUUID(),", CODE_FG)],
    [("        sequenceNumber: 0, isTruncated: false", CODE_FG)],
    [("      }],", CODE_FG)],
    [("      activityMetadata: { activity: ", CODE_FG), ('"uploadText"', CODE_STR), (" },", CODE_FG)],
    [("      protectedAppMetadata: {", CODE_FG)],
    [('        name: ', CODE_FG), ('"Agent Warden"', CODE_STR), (",", CODE_FG)],
    [("        applicationLocation: {", CODE_FG)],
    [("          value: ", CODE_FG), ("CLIENT_ID", CODE_KW), ("  ", CODE_FG), ("// d94c93dd-...", CODE_CMT)],
    [("        }", CODE_FG)],
    [("      }", CODE_FG)],
    [("    }", CODE_FG)],
    [("  };", CODE_FG)],
    "",
    [("  ", CODE_FG), ("const ", CODE_KW), ("resp = ", CODE_FG), ("await ", CODE_KW), ("fetch(GRAPH_URL, {", CODE_FG)],
    [("    method: ", CODE_FG), ('"POST"', CODE_STR), (",", CODE_FG)],
    [("    headers: { Authorization: `Bearer ${token}` },", CODE_FG)],
    [("    body: JSON.stringify(body)", CODE_FG)],
    [("  });", CODE_FG)],
    "",
    [("  ", CODE_FG), ("const ", CODE_KW), ("data = ", CODE_FG), ("await ", CODE_KW), ("resp.json();", CODE_FG)],
    [("  ", CODE_FG), ("const ", CODE_KW), ("actions = data.policyActions ?? [];", CODE_FG)],
    [("  ", CODE_FG), ("const ", CODE_KW), ("blocked = actions.some(", CODE_FG)],
    [('    a => a.action === ', CODE_FG), ('"restrictAccess"', CODE_STR)],
    [("  );", CODE_FG)],
    [("  ", CODE_FG), ("return ", CODE_KW), ("{ allowed: !blocked, actions, httpStatus };", CODE_FG)],
    [("}", CODE_FG)],
], font_size=9)

# Test runner + scenario definition
add_code_box(slide, 6.8, 1.5, 6.0, 5.5, [
    [("// Test scenario type definition", CODE_CMT)],
    [("interface ", CODE_KW), ("Scenario", CODE_FG), (" {", CODE_FG)],
    [("  name: ", CODE_FG), ("string", CODE_STR), (";", CODE_FG)],
    [("  text: ", CODE_FG), ("string", CODE_STR), (";    ", CODE_FG), ("// PII content to scan", CODE_CMT)],
    [("  expected: ", CODE_FG), ('"block" | "allow"', CODE_STR), (";", CODE_FG)],
    [("  notes: ", CODE_FG), ("string", CODE_STR), (";", CODE_FG)],
    [("}", CODE_FG)],
    "",
    [("// Test runner — iterate scenarios", CODE_CMT)],
    [("for ", CODE_KW), ("(", CODE_FG), ("const ", CODE_KW), ("scenario ", CODE_FG), ("of ", CODE_KW), ("scenarios) {", CODE_FG)],
    [("  ", CODE_FG), ("const ", CODE_KW), ("result = ", CODE_FG), ("await ", CODE_KW), ("scanText(scenario.text);", CODE_FG)],
    [("  ", CODE_FG), ("const ", CODE_KW), ('actual = result.allowed ? "allow" : "block";', CODE_FG)],
    [("  ", CODE_FG), ("const ", CODE_KW), ("pass = actual === scenario.expected;", CODE_FG)],
    [("  console.log(`${scenario.name}: ${actual} ${", CODE_FG)],
    [('    pass ? "✓" : "✗ MISMATCH"', CODE_STR), ("}`);", CODE_FG)],
    [("}", CODE_FG)],
    "",
    [("// Run command:", CODE_CMT)],
    [("// export PURVIEW_DLP_TENANT_ID=dab94ed2-...", CODE_CMT)],
    [("// export PURVIEW_DLP_USER_ID=21bbd518-...", CODE_CMT)],
    [("// export PURVIEW_DLP_CLIENT_SECRET=$(kubectl ...) ", CODE_CMT)],
    [("// node --experimental-strip-types \\", CODE_CMT)],
    [("//   test/test-dlp-scenarios-valid.ts", CODE_CMT)],
    "",
    [("// GRAPH_URL:", CODE_CMT)],
    [("// https://graph.microsoft.com/v1.0/users/", CODE_CMT)],
    [("//   {userId}/dataSecurityAndGovernance/", CODE_CMT)],
    [("//   processContent", CODE_CMT)],
], font_size=9)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Full Test Case: Scenario 1 & 2
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_stripe_rect(slide, 0, 0, 13.333, 0.08, MID_BLUE)

add_textbox(slide, 0.5, 0.3, 12, 0.7,
            "Full Test Cases — Scenarios 1 & 2",
            font_size=28, bold=True, color=DARK_BLUE)

# Scenario 1
add_textbox(slide, 0.5, 1.1, 6.0, 0.4,
            "Scenario 1: China PII (ChatGPT text input)",
            font_size=16, bold=True, color=ACCENT_PURPLE)

add_code_box(slide, 0.5, 1.6, 6.0, 2.5, [
    [("// Input text — Chinese user PII", CODE_CMT)],
    [("姓名: 王伟 (Wáng Wěi)", CODE_FG)],
    [("身份证号: 11010519900307783X", CODE_FG)],
    [("Passport number: G12345678", CODE_FG)],
    [("Credit card number: ", CODE_FG), ("6225887632109878", CODE_STR), ("  ", CODE_FG), ("← UnionPay, Luhn ✓", CODE_CMT)],
    [("地址: 北京市朝阳区建国路88号", CODE_FG)],
    "",
    [("Expected: ", CODE_FG), ("BLOCK", CODE_KW), ("   Actual: ", CODE_FG), ("BLOCK ✓", CODE_STR)],
    [("SIT matched: Credit Card Number", CODE_CMT)],
    [("HTTP 200 | restrictAccess → block", CODE_CMT)],
], font_size=10)

# Scenario 2
add_textbox(slide, 0.5, 4.3, 6.0, 0.4,
            "Scenario 2: US PII (Gemini file attachment)",
            font_size=16, bold=True, color=ACCENT_PURPLE)

add_code_box(slide, 0.5, 4.8, 6.0, 2.5, [
    [("// Input text — US user PII", CODE_CMT)],
    [("Name: Michael Anderson", CODE_FG)],
    [("SSN: ", CODE_FG), ("234-67-8901", CODE_STR), ("                  ", CODE_FG), ("← Valid format ✓", CODE_CMT)],
    [("Passport number: 285194736", CODE_FG)],
    [("Credit card number: ", CODE_FG), ("4242424242424242", CODE_STR), ("  ", CODE_FG), ("← Visa, Luhn ✓", CODE_CMT)],
    [("Address: 456 Maple Ave, Los Angeles, CA 90001", CODE_FG)],
    "",
    [("Expected: ", CODE_FG), ("BLOCK", CODE_KW), ("   Actual: ", CODE_FG), ("BLOCK ✓", CODE_STR)],
    [("SITs: Credit Card + SSN + U.S./U.K. Passport", CODE_CMT)],
    [("HTTP 200 | restrictAccess → block", CODE_CMT)],
], font_size=10)

# Response JSON
add_textbox(slide, 6.8, 1.1, 6.0, 0.4,
            "API Response (both scenarios identical)",
            font_size=16, bold=True, color=DARK_BLUE)

add_code_box(slide, 6.8, 1.6, 6.0, 2.5, [
    [("{", CODE_FG)],
    [('  "policyActions"', CODE_STR), (": [", CODE_FG)],
    [("    {", CODE_FG)],
    [('      "@odata.type"', CODE_STR), (":", CODE_FG)],
    [('        "#microsoft.graph.restrictAccessAction"', CODE_STR), (",", CODE_FG)],
    [('      "action"', CODE_STR), (": ", CODE_FG), ('"restrictAccess"', CODE_STR), (",", CODE_FG)],
    [('      "restrictionAction"', CODE_STR), (": ", CODE_FG), ('"block"', CODE_STR)],
    [("    }", CODE_FG)],
    [("  ]", CODE_FG)],
    [("}", CODE_FG)],
], font_size=10)

# Checksum validation notes
tf = add_rich_textbox(slide, 6.8, 4.3, 6.0, 3.0)
add_para(tf, "Checksum Validation Notes", font_size=16, bold=True, color=DARK_BLUE)
add_para(tf, "", font_size=6)
add_para(tf, "All credit card numbers pass the Luhn algorithm — required for the Credit Card Number SIT to trigger.", font_size=12, color=BLACK)
add_para(tf, "", font_size=4)
add_para(tf, "• UnionPay 6225... — BIN range 622126-622925", font_size=11, color=GRAY)
add_para(tf, "• Visa 4242... — standard test card", font_size=11, color=GRAY)
add_para(tf, "• JCB 3528... — BIN range 3528-3589", font_size=11, color=GRAY)
add_para(tf, "• MasterCard 5555... — BIN range 5100-5599", font_size=11, color=GRAY)
add_para(tf, "• MC 2221... — BIN range 2221-2720", font_size=11, color=GRAY)
add_para(tf, "", font_size=4)
add_para(tf, "Keywords must be in the correct language for corroborative evidence (e.g., クレジットカード for Japanese).", font_size=11, color=GRAY)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Full Test Case: Scenario 3 & 4
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_stripe_rect(slide, 0, 0, 13.333, 0.08, MID_BLUE)

add_textbox(slide, 0.5, 0.3, 12, 0.7,
            "Full Test Cases — Scenarios 3 & 4",
            font_size=28, bold=True, color=DARK_BLUE)

# Scenario 3
add_textbox(slide, 0.5, 1.1, 6.0, 0.4,
            "Scenario 3: Japan PII (ZIP file attachment)",
            font_size=16, bold=True, color=ACCENT_PURPLE)

add_code_box(slide, 0.5, 1.6, 6.0, 2.5, [
    [("// Input text — Japanese user PII (native keywords)", CODE_CMT)],
    [("氏名: 佐藤 花子 (Sato Hanako)", CODE_FG)],
    [("マイナンバー: ", CODE_FG), ("123456789018", CODE_STR), ("        ", CODE_FG), ("← Valid check digit ✓", CODE_CMT)],
    [("パスポート番号: TR1234567", CODE_FG)],
    [("クレジットカード番号: ", CODE_FG), ("3528234567890126", CODE_STR), (" ", CODE_FG), ("← JCB, Luhn ✓", CODE_CMT)],
    [("住所: 東京都渋谷区神宮前1-2-3", CODE_FG)],
    "",
    [("Expected: ", CODE_FG), ("BLOCK", CODE_KW), ("   Actual: ", CODE_FG), ("BLOCK ✓", CODE_STR)],
    [("SITs: Credit Card + My Number Personal + Passport", CODE_CMT)],
    [("HTTP 200 | restrictAccess → block", CODE_CMT)],
], font_size=10)

# Scenario 4
add_textbox(slide, 0.5, 4.3, 6.0, 0.4,
            "Scenario 4: France PII (Claude.ai interaction)",
            font_size=16, bold=True, color=ACCENT_PURPLE)

add_code_box(slide, 0.5, 4.8, 6.0, 2.5, [
    [("// Input text — French user PII (French keywords)", CODE_CMT)],
    [("Nom: Marie Dupont-Lefèvre", CODE_FG)],
    [("Numéro de sécurité sociale: ", CODE_FG), ("185067283009381", CODE_STR), (" ", CODE_FG), ("← INSÉE key=81 ✓", CODE_CMT)],
    [("Numéro de passeport: 10AB12345", CODE_FG)],
    [("Numéro de carte de crédit: ", CODE_FG), ("5555555555554444", CODE_STR), (" ", CODE_FG), ("← MC, Luhn ✓", CODE_CMT)],
    [("Adresse: 15 Rue de Rivoli, 75001 Paris", CODE_FG)],
    "",
    [("Expected: ", CODE_FG), ("BLOCK", CODE_KW), ("   Actual: ", CODE_FG), ("BLOCK ✓", CODE_STR)],
    [("SITs: Credit Card + France Passport", CODE_CMT)],
    [("HTTP 200 | restrictAccess → block", CODE_CMT)],
], font_size=10)

# Scenario correction notes
tf = add_rich_textbox(slide, 6.8, 1.1, 6.0, 3.0)
add_para(tf, "Data Corrections Applied", font_size=16, bold=True, color=DARK_BLUE)
add_para(tf, "", font_size=6)
add_para(tf, "Scenario 3 — Japan:", font_size=14, bold=True, color=ACCENT_PURPLE)
add_para(tf, "• JCB: 3528234567890123 → 3528234567890126 (Luhn fix)", font_size=11, color=BLACK)
add_para(tf, "• My Number: 123456789012 → 123456789018 (check digit fix)", font_size=11, color=BLACK)
add_para(tf, "• Added Japanese keywords: マイナンバー, パスポート番号, クレジットカード番号", font_size=11, color=BLACK)
add_para(tf, "", font_size=6)
add_para(tf, "Scenario 4 — France:", font_size=14, bold=True, color=ACCENT_PURPLE)
add_para(tf, "• INSÉE: 1850672830093 → 185067283009381 (added 2-digit key: 81)", font_size=11, color=BLACK)
add_para(tf, "• Added French keywords: Numéro de passeport, Numéro de carte de crédit", font_size=11, color=BLACK)

tf = add_rich_textbox(slide, 6.8, 4.3, 6.0, 3.0)
add_para(tf, "Why Corroborative Keywords Matter", font_size=16, bold=True, color=DARK_BLUE)
add_para(tf, "", font_size=6)
add_para(tf, "Purview SITs use a two-part detection model:", font_size=12, color=BLACK)
add_para(tf, "1. Primary pattern — regex/checksum match (e.g., Luhn algorithm for credit cards)", font_size=11, color=BLACK)
add_para(tf, '2. Corroborative evidence — nearby keywords (e.g., "credit card", "クレジットカード")', font_size=11, color=BLACK)
add_para(tf, "", font_size=4)
add_para(tf, "Without corroborative keywords, confidence stays Low even if the number is valid. The DLP rule requires Medium (85+) confidence.", font_size=11, color=GRAY)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Full Test Case: Scenario 5
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_stripe_rect(slide, 0, 0, 13.333, 0.08, MID_BLUE)

add_textbox(slide, 0.5, 0.3, 12, 0.7,
            "Full Test Case — Scenario 5 & Terminal Output",
            font_size=28, bold=True, color=DARK_BLUE)

# Scenario 5
add_textbox(slide, 0.5, 1.1, 6.0, 0.4,
            "Scenario 5: Russia PII (photo text extraction)",
            font_size=16, bold=True, color=ACCENT_PURPLE)

add_code_box(slide, 0.5, 1.6, 6.0, 2.7, [
    [("// Input text — Russian user PII (Cyrillic keywords)", CODE_CMT)],
    [("ФИО: Иванов Алексей Петрович", CODE_FG)],
    [("ИНН: ", CODE_FG), ("770123456703", CODE_STR), ("                 ", CODE_FG), ("← Valid check digits ✓", CODE_CMT)],
    [("Номер паспорта: 45 1234567", CODE_FG)],
    [("Номер кредитной карты: ", CODE_FG), ("2221001234567896", CODE_STR), (" ", CODE_FG), ("← MC 2221, Luhn ✓", CODE_CMT)],
    [("Адрес: ул. Тверская, д. 15, Москва, 125009", CODE_FG)],
    "",
    [("Expected: ", CODE_FG), ("BLOCK", CODE_KW), ("   Actual: ", CODE_FG), ("BLOCK ✓", CODE_STR)],
    [("SITs: Credit Card + Russian Passport + Taxpayer ID", CODE_CMT)],
    [("HTTP 200 | restrictAccess → block", CODE_CMT)],
    [("Note: Mir (2200) not recognized by Purview — used MC 2221 prefix", CODE_CMT)],
], font_size=10)

# Terminal output
add_textbox(slide, 0.5, 4.6, 12.3, 0.4,
            "Terminal Output — Full Test Run (April 15, 2026)",
            font_size=16, bold=True, color=DARK_BLUE)

add_code_box(slide, 0.5, 5.1, 12.3, 2.2, [
    [("$ ", CODE_CMT), ("node --experimental-strip-types test/test-dlp-scenarios-valid.ts", CODE_FG)],
    [("✓ Token acquired successfully", CODE_STR)],
    "",
    [("━━━ Scenario 1: China PII (valid Luhn)            HTTP: 200  Result: ", CODE_FG), ("BLOCK ✓", CODE_STR)],
    [("━━━ Scenario 2: US PII (valid Luhn + SSN)         HTTP: 200  Result: ", CODE_FG), ("BLOCK ✓", CODE_STR)],
    [("━━━ Scenario 3: Japan PII (valid Luhn + My Number) HTTP: 200  Result: ", CODE_FG), ("BLOCK ✓", CODE_STR)],
    [("━━━ Scenario 4: France PII (valid Luhn + INSÉE)   HTTP: 200  Result: ", CODE_FG), ("BLOCK ✓", CODE_STR)],
    [("━━━ Scenario 5: Russia PII (valid Luhn + INN)     HTTP: 200  Result: ", CODE_FG), ("BLOCK ✓", CODE_STR)],
    "",
    [("Total: 5  Passed: ", CODE_FG), ("5", CODE_STR), ("  Failed: ", CODE_FG), ("0", CODE_STR)],
], font_size=10)

# Data correction note
tf = add_rich_textbox(slide, 6.8, 1.1, 6.0, 3.2)
add_para(tf, "Scenario 5 — Data Corrections", font_size=16, bold=True, color=DARK_BLUE)
add_para(tf, "", font_size=6)
add_para(tf, "Credit Card:", font_size=13, bold=True, color=ACCENT_PURPLE)
add_para(tf, "• Original: Mir 2200123456789012 — Luhn valid but Purview doesn't recognize Mir BIN range (2200-2204)", font_size=11, color=BLACK)
add_para(tf, "• Fixed: MC 2221001234567896 — BIN 2221-2720 is in MasterCard range, recognized by Purview", font_size=11, color=BLACK)
add_para(tf, "", font_size=4)
add_para(tf, "INN (Taxpayer ID):", font_size=13, bold=True, color=ACCENT_PURPLE)
add_para(tf, "• Original: 770123456789 — invalid check digits", font_size=11, color=BLACK)
add_para(tf, "• Fixed: 770123456703 — valid per Russian INN algorithm", font_size=11, color=BLACK)
add_para(tf, "", font_size=4)
add_para(tf, "Russian keywords:", font_size=13, bold=True, color=ACCENT_PURPLE)
add_para(tf, "• ИНН (taxpayer ID), Номер паспорта (passport number), Номер кредитной карты (credit card number)", font_size=11, color=BLACK)


# ── Save ──
output_path = "/Users/hifrankc/workspace/agent-warden/docs/Agent-Warden-DLP-Test-Results.pptx"
prs.save(output_path)
print(f"✓ PowerPoint saved to: {output_path}")
